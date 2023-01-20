import os
from pptx import Presentation

def remove_vt_from_pptx(file_path):
    # Open the PPTX file
    prs = Presentation(file_path)

    # Create a log file to store the removed VTs
    log_file = open("removed_vts.txt", "w")

    # Iterate through each slide in the PPTX file
    for slide in prs.slides:
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has a text frame
            if hasattr(shape, "text_frame"):
                # Iterate through each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Iterate through each run in the paragraph
                    for run in paragraph.runs:
                        # Get the text of the run
                        text = run.text
                        # Replace all VTs with spaces if there is no space before or after
                        new_text = text.replace("\v", " ") if " \v " not in text else text
                        # Delete all VTs if there is a space before or after
                        new_text = new_text.replace(" \v ", "")
                        # Update the text of the run
                        run.text = new_text
                        if text != new_text:
                            log_file.write("Removed in slide {}, text {}".format(slide.slide_id, text))
    log_file.close()
    # Save the changes to the PPTX file
    prs.save(file_path)
    
remove_vt_from_pptx('SampleFile.pptx')